import os
import platform
import time
import json
import re
import logging
import shutil
import textwrap
import sys
import glob
import html
import csv
import io
import hashlib
import tempfile
import gc
from collections import OrderedDict
import threading

try:
    import cv2
except ImportError:  # pragma: no cover - optional dependency in lightweight test envs
    cv2 = None

try:
    from PIL import Image
except ImportError:  # pragma: no cover - optional dependency in lightweight test envs
    Image = None

try:
    import fitz
except ImportError:  # pragma: no cover - optional dependency in lightweight test envs
    fitz = None

try:
    import docx
except ImportError:  # pragma: no cover - optional dependency in lightweight test envs
    docx = None

try:
    from fpdf import FPDF
except ImportError:  # pragma: no cover - optional dependency in lightweight test envs
    FPDF = None

try:
    from google import genai
except ImportError:  # pragma: no cover - optional dependency in lightweight test envs
    genai = None

try:
    import openpyxl
except ImportError:  # pragma: no cover - optional dependency in lightweight test envs
    openpyxl = None

try:
    from ebooklib import epub
except ImportError:  # pragma: no cover - optional dependency in lightweight test envs
    epub = None
from chronicle_app.services.app_files import (
    emit_launch_continuity as shared_emit_launch_continuity,
    resolve_runtime_crash_log_path as shared_resolve_runtime_crash_log_path,
)
from chronicle_app.services.prompting import (
    build_prompt as build_shared_prompt,
    enforce_archival_heading_structure as shared_enforce_archival_heading_structure,
)
from chronicle_app.config import (
    PROFILE_CHOICES,
    PROFILE_PRESETS,
    RTL_LANGUAGE_CODES,
    TRANSLATION_TARGETS,
)
from chronicle_core import (
    apply_newspaper_html_safety_fallback as core_apply_newspaper_html_safety_fallback,
    apply_modern_punctuation as core_apply_modern_punctuation,
    apply_modern_currency as core_apply_modern_currency,
    apply_expanded_abbreviations as core_apply_expanded_abbreviations,
    apply_handwriting_audit_flag as core_apply_handwriting_audit_flag,
    build_tabular_html_fragment as core_build_tabular_html_fragment,
    build_newspaper_safety_notice as core_build_newspaper_safety_notice,
    clean_text_artifacts as core_clean_text_artifacts,
    csv_to_accessible_text as core_csv_to_accessible_text,
    get_newspaper_profile_rules as core_get_newspaper_profile_rules,
    parse_csv_rows as core_parse_csv_rows,
    recover_newspaper_header_citation as core_recover_newspaper_header_citation,
    recover_source_attribution_footer as core_recover_source_attribution_footer,
    sanitize_latin1 as core_sanitize_latin1,
    sanitize_model_output as core_sanitize_model_output,
    normalize_streamed_html_document as core_normalize_streamed_html_document,
    write_header as core_write_header,
    write_footer as core_write_footer,
)
from legacy_extract_core import (
    process_pdf_gemini as legacy_process_pdf_gemini,
)

# Dynamically link Homebrew's Python site-packages for Mac users
brew_paths = glob.glob("/opt/homebrew/lib/python3.*/site-packages") + glob.glob("/usr/local/lib/python3.*/site-packages")
for path in brew_paths:
    if os.path.exists(path) and path not in sys.path:
        sys.path.append(path)


APP_NAME = "Chronicle"

def _get_crash_log_path():
    return shared_resolve_runtime_crash_log_path(app_name=APP_NAME)

if getattr(sys, 'frozen', False):
    try:
        crash_log_path = _get_crash_log_path()
        sys.stdout = open(crash_log_path, 'a', buffering=1, encoding='utf-8', errors='replace')
        sys.stderr = sys.stdout
    except Exception:
        pass

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
INPUT_DIR = os.path.join(SCRIPT_DIR, "input_files")
BATCH_INPUT_DIR = os.path.join(SCRIPT_DIR, "Input_Scans")

def _get_app_data_dir():
    if platform.system() == "Windows":
        base = os.environ.get("APPDATA") or os.path.expanduser("~")
        return os.path.join(base, APP_NAME)
    if platform.system() == "Darwin":
        return os.path.join(os.path.expanduser("~/Library/Application Support"), APP_NAME)
    base = os.environ.get("XDG_CONFIG_HOME") or os.path.expanduser("~/.config")
    return os.path.join(base, APP_NAME)

APP_DATA_DIR = _get_app_data_dir()
os.makedirs(APP_DATA_DIR, exist_ok=True)
CONFIG_FILE = os.path.join(APP_DATA_DIR, "user_config.json")

PDF_CHUNK_PAGES = 5
TEXT_CHUNK_CHARS = 15000
HEARTBEAT_TIMEOUT_SEC = 480
API_POLITE_MODE = True
API_MIN_REQUEST_INTERVAL_SEC = 1.5 if API_POLITE_MODE else 1.0
API_MAX_CONCURRENT_REQUESTS = 1
API_MAX_PENDING_REQUESTS = 24
API_REQUEST_QUEUE_POLL_SEC = 0.25
API_CACHE_MAX_ENTRIES = 3000
TEXT_BATCH_TARGET_CHARS = 24000
API_CONNECTION_FAILURE_DELAY_SEC = 20.0
GEMINI_UPLOAD_POLL_SEC = 2.5
GEMINI_UPLOAD_MAX_WAIT_SEC = 180.0
SUPPORTED_EXTENSIONS = ['.pdf', '.docx', '.txt', '.md', '.rtf', '.csv', '.js', '.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp', '.xlsx', '.pptx']
PROTECTED_INPUT_DIR_BASENAMES = {"checking documents"}
TRANSLATION_TARGETS = [
    ("English", "en"),
    ("Spanish", "es"),
    ("French", "fr"),
    ("German", "de"),
    ("Italian", "it"),
    ("Portuguese", "pt"),
]
RTL_LANGUAGE_CODES = {"ar", "fa", "he", "ur", "yi", "ps", "sd", "ug"}
_api_request_lock = threading.Lock()
_api_last_request_ts = 0.0
_api_request_semaphore = threading.Semaphore(API_MAX_CONCURRENT_REQUESTS)
_api_queue_lock = threading.Lock()
_api_pending_requests = 0
_chunk_cache = OrderedDict()
_chunk_cache_lock = threading.Lock()

def _secure_chmod(path, mode=0o600):
    if os.name == "nt":
        return
    try:
        os.chmod(path, mode)
    except Exception:
        pass

def migrate_legacy_cli_key_files():
    for vendor in ("gemini", "anthropic", "openai"):
        legacy = os.path.join(SCRIPT_DIR, f"api_key_{vendor}.txt")
        target = os.path.join(APP_DATA_DIR, f"api_key_{vendor}.txt")
        if os.path.exists(legacy) and not os.path.exists(target):
            try:
                shutil.copy2(legacy, target)
                _secure_chmod(target, 0o600)
            except Exception:
                pass
        if os.path.exists(legacy) and os.path.exists(target):
            try:
                os.remove(legacy)
            except Exception:
                try:
                    with open(legacy, "w", encoding="utf-8") as fh:
                        fh.write("")
                except Exception:
                    pass

def get_translation_target(config):
    target = str(config.get("translate_target", "English")).strip()
    lookup = {name.lower(): (name, code) for name, code in TRANSLATION_TARGETS}
    return lookup.get(target.lower(), ("English", "en"))

def get_output_lang_code(config):
    if config.get("translate_mode", "none") == "none":
        return "und"
    _, code = get_translation_target(config)
    return code

def get_output_text_direction(config):
    if config.get("translate_mode", "none") == "none":
        return "auto"
    return "rtl" if get_output_lang_code(config) in RTL_LANGUAGE_CODES else "ltr"

def is_path_within_protected_input_dirs(path):
    try:
        target = os.path.normcase(os.path.abspath(path))
    except Exception:
        return False
    for dirname in PROTECTED_INPUT_DIR_BASENAMES:
        root = os.path.normcase(os.path.abspath(os.path.join(SCRIPT_DIR, dirname)))
        try:
            if os.path.commonpath([target, root]) == root:
                return True
        except Exception:
            continue
    return False

def get_api_key(model_name):
    vendor = "claude" if "claude" in model_name else "openai" if "gpt" in model_name else "gemini"
    candidate_files = [os.path.join(APP_DATA_DIR, f"api_key_{vendor}.txt")]
    if vendor == "claude":
        candidate_files.append(os.path.join(APP_DATA_DIR, "api_key_anthropic.txt"))
    for key_file in candidate_files:
        if os.path.exists(key_file):
            with open(key_file, "r") as f:
                key = f.read().strip()
                if key:
                    return key
    print(f"\nFIRST TIME SETUP - {vendor.upper()}")
    new_key = input(f"Please paste your {vendor.upper()} API Key here and press Enter: ").strip()
    if new_key:
        key_file = candidate_files[0]
        with open(key_file, "w") as f: f.write(new_key)
        _secure_chmod(key_file, 0o600)
        return new_key
    exit()

def ask_menu(title, options, option_map, default_key=''):
    print(f"\n{title}")
    for opt in options: print(opt)
    while True:
        choice = input("Select option: ").strip()
        if choice in option_map: return option_map[choice]
        if choice == '' and default_key in option_map: return option_map[default_key]
        print("Invalid choice.")

def ask_bool(title, options, true_choice='2', default_choice='0'):
    print(f"\n{title}")
    for opt in options: print(opt)
    while True:
        choice = input("Select option: ").strip()
        if choice == '':
            choice = default_choice
        if choice == true_choice: return True
        if choice in ['0', '1']: return False
        print("Invalid choice.")

def get_user_preferences():
    config = {}
    fmt_map = {'1': ('html', 'output_html'), '2': ('txt', 'output_txt'), '3': ('docx', 'output_docx'), '4': ('md', 'output_md'), '5': ('pdf', 'output_pdf'), '6': ('json', 'output_json'), '7': ('csv', 'output_csv'), '8': ('epub', 'output_epub')}
    fmt, out_folder = ask_menu("MENU 1: OUTPUT FORMAT", ["1. HTML", "2. TXT", "3. DOCX", "4. MD", "5. PDF", "6. JSON", "7. CSV", "8. EPUB"], fmt_map)
    config['format_type'] = fmt
    config['output_dir'] = os.path.join(SCRIPT_DIR, out_folder)
    model_options = [
        "1. Gemini 2.5 Flash",
        "2. Gemini 2.5 Pro (Deep Scan)",
        "3. Claude Sonnet 4",
        "4. GPT-4o",
    ]
    model_map = {'1':'gemini-2.5-flash', '2':'gemini-2.5-pro', '3':'claude-sonnet-4-20250514', '4':'gpt-4o', '':'gemini-2.5-flash'}
    config['model_name'] = ask_menu("MENU 2: AI ENGINE", model_options, model_map, '')

    profile_options = [f"{idx}. {label}" for idx, (_, label) in enumerate(PROFILE_CHOICES, start=1)]
    profile_map = {str(idx): key for idx, (key, _) in enumerate(PROFILE_CHOICES, start=1)}
    config['doc_profile'] = ask_menu("MENU 3: DOCUMENT PROFILE", profile_options, profile_map, '1')
    shared_preset = PROFILE_PRESETS.get(config['doc_profile'], PROFILE_PRESETS['standard'])
    recommended_model = str(shared_preset.get('model_name', 'gemini-2.5-flash'))
    model_default_map = {
        'gemini-2.5-flash': '1',
        'gemini-2.5-pro': '2',
        'claude-sonnet-4-20250514': '3',
        'gpt-4o': '4',
    }
    preset = {
        'model': model_default_map.get(recommended_model, '1'),
        'translate': '0' if shared_preset.get('translate_mode', 'none') == 'none' else ('2' if shared_preset.get('translate_mode') == 'both' else '3'),
        'punct': bool(shared_preset.get('modernize_punctuation', False)),
        'unit': bool(shared_preset.get('unit_conversion', False)),
        'merge': bool(shared_preset.get('merge_files', False)),
        'image': bool(shared_preset.get('image_descriptions', True)),
        'large': bool(shared_preset.get('large_print', False)),
        'abbrev': bool(shared_preset.get('abbrev_expansion', False)),
    }
    config['model_name'] = ask_menu("MENU 3B: RECOMMENDED ENGINE (BY PRESET)", model_options, model_map, preset['model'])
    
    config['translate_mode'] = ask_menu("MENU 4: LANGUAGE TRANSLATION", ["0. Disable translation (keep original language)", "2. Keep original + translation in brackets", "3. Translate only"], {'0':'none', '2':'both', '3':'english_only', '':'none'}, preset['translate'])
    if config['translate_mode'] != 'none':
        target_map = {'1': 'English', '2': 'Spanish', '3': 'French', '4': 'German', '5': 'Italian', '6': 'Portuguese', '': 'English'}
        config['translate_target'] = ask_menu("MENU 4B: TRANSLATION TARGET", ["1. English", "2. Spanish", "3. French", "4. German", "5. Italian", "6. Portuguese"], target_map, '')
    else:
        config['translate_target'] = 'English'
    config['modernize_punctuation'] = ask_bool("MENU 5: PUNCTUATION", ["0. Strictly preserve original punctuation", "2. Update to modern punctuation"], default_choice='2' if preset['punct'] else '0')
    config['unit_conversion'] = ask_bool("MENU 6: CURRENCY & MEASUREMENTS", ["0. Leave old currency/measurements untouched", "2. Keep old values and add modern equivalence in brackets"], default_choice='2' if preset['unit'] else '0')
    config['merge_files'] = ask_bool("MENU 7: FILE HANDLING", ["0. Individual", "2. Merge files"], default_choice='2' if preset['merge'] else '0')
    config['collision_mode'] = ask_menu("MENU 8: FILE COLLISIONS", ["1. Skip (Default)", "2. Overwrite", "3. Auto-Number"], {'1':'skip', '2':'overwrite', '3':'auto'}, '1')
    config['image_descriptions'] = ask_menu("MENU 9: VISUALS", ["0. Enable", "1. Disable"], {'0':True, '1':False, '':True}, '0' if preset['image'] else '1')
    config['large_print'] = ask_bool("MENU 10: LARGE PRINT PDF", ["0. Standard 11pt", "2. Enable 18pt High-Contrast"], default_choice='2' if preset['large'] else '0') if fmt == 'pdf' else False
    config['batch_mode'] = ask_menu("MENU 11: BATCH SCAN", ["0. Standard", "1. Recursive (Keep)", "2. Recursive (Delete)"], {'0':'flat', '1':'recursive_keep', '2':'recursive_delete', '':'flat'}, '')
    config['abbrev_expansion'] = ask_bool("MENU 12: ABBREVIATIONS", ["0. Skip", "2. Expand Contextually (e.g., Bn [Battalion])"], default_choice='2' if preset['abbrev'] else '0')
    config['academic_mode'] = config['doc_profile'] == 'academic'
    return config

def get_prompt(config):
    return build_shared_prompt(
        config,
        translation_targets=TRANSLATION_TARGETS,
        rtl_language_codes=RTL_LANGUAGE_CODES,
        format_type=config.get("format_type", "html"),
    )
def enhance_image_for_microtext(file_path):
    try:
        img = cv2.imread(file_path)
        if img is None:
            return file_path
        width, height = int(img.shape[1] * 2), int(img.shape[0] * 2)
        resized = cv2.resize(img, (width, height), interpolation=cv2.INTER_LANCZOS4)
        gray = cv2.cvtColor(resized, cv2.COLOR_BGR2GRAY)
        # Improve faint handwriting and soft scans without inventing content.
        denoised = cv2.fastNlMeansDenoising(gray, None, 8, 7, 21)
        contrast = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8)).apply(denoised)
        processed = cv2.addWeighted(contrast, 1.22, cv2.GaussianBlur(contrast, (0, 0), 1.1), -0.22, 0)
        safe_stem = re.sub(r"[^0-9A-Za-z._-]+", "_", os.path.basename(file_path))
        temp_path = os.path.join(
            tempfile.gettempdir(),
            f"chronicle_enhanced_{os.getpid()}_{safe_stem}.png",
        )
        if not cv2.imwrite(temp_path, processed):
            return file_path
        return temp_path
    except Exception as e:
        print(f"Warning: Could not enhance image ({e}). Using original.")
        return file_path

def clean_text_artifacts(text):
    return core_clean_text_artifacts(text)

def csv_to_accessible_text(raw_text, max_rows=None, max_cell_chars=None):
    return core_csv_to_accessible_text(raw_text, max_rows=max_rows, max_cell_chars=max_cell_chars)

def sanitize_latin1(text):
    return core_sanitize_latin1(text)

def sanitize_model_output(text_content, format_type, doc_profile=None, preserve_original_page_numbers=False):
    return core_sanitize_model_output(
        text_content,
        format_type,
        doc_profile,
        preserve_original_page_numbers,
    )

def apply_modern_punctuation(text_content):
    return core_apply_modern_punctuation(text_content)

def apply_modern_currency(text_content):
    return core_apply_modern_currency(text_content)


def apply_expanded_abbreviations(text_content):
    return core_apply_expanded_abbreviations(text_content)

def normalize_streamed_html_document(full_html):
    return core_normalize_streamed_html_document(full_html)


def enforce_archival_heading_structure(content, fmt, doc_profile):
    return shared_enforce_archival_heading_structure(content, fmt, doc_profile)

def apply_handwriting_audit_flag(text_content, format_type, doc_profile, whole_document=False):
    return core_apply_handwriting_audit_flag(text_content, format_type, doc_profile, whole_document=whole_document)

def apply_newspaper_html_safety_fallback(text_content, format_type, doc_profile, max_chars=14000):
    return core_apply_newspaper_html_safety_fallback(text_content, format_type, doc_profile, max_chars=max_chars)

def recover_source_attribution_footer(text_content, format_type, doc_profile, source_path=None):
    return core_recover_source_attribution_footer(
        text_content, format_type, doc_profile, source_path=source_path
    )


def recover_newspaper_header_citation(text_content, format_type, doc_profile, source_path=None):
    return core_recover_newspaper_header_citation(
        text_content, format_type, doc_profile, source_path=source_path
    )

def write_header(file_obj, title, format_type, lang_code='en', text_dir='ltr'):
    core_write_header(file_obj, title, format_type, lang_code=lang_code, text_dir=text_dir)

def write_footer(file_obj, format_type):
    core_write_footer(file_obj, format_type)

def _probe_pdf_text_layer(pdf_path, sample_pages=3, min_chars_per_page=200):
    try:
        doc = fitz.open(pdf_path)
    except Exception:
        return False
    try:
        total_pages = len(doc)
        if total_pages <= 0:
            return False
        pages_to_check = min(max(1, sample_pages), total_pages)
        total_chars = 0
        for page_index in range(pages_to_check):
            try:
                extracted = doc.load_page(page_index).get_text("text") or ""
            except Exception:
                extracted = ""
            total_chars += len(extracted.strip())
        average_chars = total_chars / pages_to_check
        return average_chars > min_chars_per_page
    finally:
        doc.close()

def _render_pdf_page_to_png_bytes(doc, page_index, zoom=2.0):
    page = doc.load_page(page_index)
    matrix = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    image_bytes = pix.tobytes("png")
    pix = None
    page = None
    return image_bytes


def _write_temp_png(prefix, image_bytes):
    fd, temp_path = tempfile.mkstemp(prefix=prefix, suffix=".png")
    try:
        with os.fdopen(fd, "wb") as fh:
            fh.write(image_bytes)
    except Exception:
        try:
            os.close(fd)
        except Exception:
            pass
        if os.path.exists(temp_path):
            os.remove(temp_path)
        raise
    return temp_path


def save_as_pdf(pdf_path, text_content, large_print=False):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    font_size = 18 if large_print else 11
    pdf.set_font("Helvetica", size=font_size)
    if large_print:
        # Force FPDF HTML parser to respect large-font rendering.
        text_content = f"<div style='font-size: 24px; line-height: 1.5;'>{text_content}</div>"
    try:
        pdf.write_html(text_content)
        pdf.output(pdf_path)
    except Exception as e:
        print(f"Warning: HTML-PDF compilation failed, falling back to flat text. {e}")
        fallback = FPDF()
        fallback.add_page()
        fallback.set_auto_page_break(auto=True, margin=15)
        fallback.set_font("Helvetica", size=font_size)

        # Strip HTML and unescape entities for clean plain-text fallback.
        import html as html_lib
        clean_fallback = re.sub(r'<[^>]+>', '', text_content)
        clean_fallback = html_lib.unescape(clean_fallback)

        fallback.multi_cell(0, 10 if not large_print else 16, text=sanitize_latin1(clean_fallback))
        fallback.output(pdf_path)

def append_to_docx(docx_path, text_content):
    doc = docx.Document(docx_path) if os.path.exists(docx_path) else docx.Document()
    for line in text_content.split('\n'):
        clean_line = line.strip()
        if clean_line.startswith('### '): doc.add_heading(clean_line[4:], level=3)
        elif clean_line.startswith('## '): doc.add_heading(clean_line[3:], level=2)
        elif clean_line.startswith('# '): doc.add_heading(clean_line[2:], level=1)
        elif clean_line.startswith('-') or clean_line.startswith('*'):
            clean_text = clean_line.lstrip('-* ').strip()
            doc.add_paragraph(clean_text, style='List Bullet')
        elif clean_line != "": doc.add_paragraph(clean_line)
    doc.save(docx_path)

def save_as_json(json_path, text_content):
    text_content = text_content.strip()
    if text_content.startswith("```json"): text_content = text_content[7:-3].strip()
    try: data = json.loads(text_content)
    except: data = {"chronicle_extracted_content": text_content}
    with open(json_path, 'w', encoding='utf-8') as f: json.dump(data, f, indent=4)

def save_as_csv(csv_path, text_content):
    text_content = text_content.replace("```csv", "").replace("```", "").strip()
    with open(csv_path, 'w', encoding='utf-8') as f: f.write(text_content)

def save_as_epub(epub_path, title, text_content, lang_code='en', text_dir='ltr'):
    book = epub.EpubBook()
    book.set_identifier(f"chron_{int(time.time())}")
    book.set_title(title)
    book.set_language(lang_code)
    
    chapters = re.split(r'(<h[123].*?>.*?</h[123]>)', text_content, flags=re.IGNORECASE)
    epub_chapters, current_title, current_content, chap_idx = [], title, "", 1
    
    for segment in chapters:
        if segment.lower().startswith('<h1') or segment.lower().startswith('<h2') or segment.lower().startswith('<h3'):
            if current_content.strip():
                c = epub.EpubHtml(title=current_title, file_name=f'chap_{chap_idx}.xhtml', lang=lang_code)
                c.content = f"<h1>{current_title}</h1><div dir=\"{text_dir}\">{current_content}</div>"
                epub_chapters.append(c)
                chap_idx += 1
            current_title = re.sub(r'<[^>]+>', '', segment)
            current_content = segment 
        else:
            current_content += segment
            
    if current_content.strip():
        c = epub.EpubHtml(title=current_title, file_name=f'chap_{chap_idx}.xhtml', lang=lang_code)
        c.content = f"<h1>{current_title}</h1><div dir=\"{text_dir}\">{current_content}</div>"
        epub_chapters.append(c)

    for c in epub_chapters: book.add_item(c)
    book.spine = ['nav'] + epub_chapters
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(epub_path, book, {})

def dispatch_save(config, path, memory_list, title, clear_memory=False):
    content = "".join(memory_list)
    if not content: return
    fmt = config['format_type']
    content = sanitize_model_output(
        content,
        fmt,
        config.get('doc_profile'),
        config.get('preserve_original_page_numbers', False),
    )
    if config.get('modernize_punctuation'):
        content = apply_modern_punctuation(content)
    if config.get('unit_conversion'):
        content = apply_modern_currency(content)
    if config.get('abbrev_expansion'):
        content = apply_expanded_abbreviations(content)
    content = apply_newspaper_html_safety_fallback(content, fmt, config.get('doc_profile'))
    content = recover_newspaper_header_citation(content, fmt, config.get('doc_profile'), config.get('source_path'))
    content = recover_source_attribution_footer(content, fmt, config.get('doc_profile'), config.get('source_path'))
    content = apply_handwriting_audit_flag(content, fmt, config.get('doc_profile'))
    content = enforce_archival_heading_structure(content, fmt, config.get('doc_profile'))
    
    if fmt == 'docx': 
        append_to_docx(path, content)
        if clear_memory: memory_list.clear()
    elif fmt == 'pdf': save_as_pdf(path, content, config.get('large_print', False))
    elif fmt == 'json': save_as_json(path, content)
    elif fmt == 'csv': save_as_csv(path, content)
    elif fmt == 'epub': save_as_epub(path, title, content, get_output_lang_code(config), get_output_text_direction(config))
def pace_api_request():
    global _api_last_request_ts
    with _api_request_lock:
        now = time.time()
        wait = API_MIN_REQUEST_INTERVAL_SEC - (now - _api_last_request_ts)
        if wait > 0:
            time.sleep(wait)
        _api_last_request_ts = time.time()

def wait_for_request_slot():
    global _api_pending_requests
    while True:
        with _api_queue_lock:
            if _api_pending_requests < API_MAX_PENDING_REQUESTS:
                _api_pending_requests += 1
                break
        time.sleep(API_REQUEST_QUEUE_POLL_SEC)
    _api_request_semaphore.acquire()
    with _api_queue_lock:
        _api_pending_requests = max(0, _api_pending_requests - 1)

def release_request_slot():
    _api_request_semaphore.release()

def is_connection_path_error(exc):
    err = str(exc).lower()
    return any(token in err for token in (
        "connecterror",
        "connection error",
        "temporary failure in name resolution",
        "nodename nor servname provided",
        "name or service not known",
        "getaddrinfo",
        "failed to establish a new connection",
        "connection refused",
        "network is unreachable",
    ))

def sha256_text(text):
    return hashlib.sha256(text.encode("utf-8", errors="ignore")).hexdigest()

def sha256_file(path):
    digest = hashlib.sha256()
    with open(path, "rb") as fh:
        while True:
            block = fh.read(1024 * 1024)
            if not block:
                break
            digest.update(block)
    return digest.hexdigest()

def build_request_cache_key(model_name, prompt_text, payload_kind, payload_fingerprint):
    return f"{model_name}|{payload_kind}|{sha256_text(prompt_text)}|{payload_fingerprint}"

def cache_get(cache_key):
    with _chunk_cache_lock:
        val = _chunk_cache.get(cache_key)
        if val is None:
            return None
        _chunk_cache.move_to_end(cache_key)
        return val

def cache_put(cache_key, text):
    if not text:
        return
    with _chunk_cache_lock:
        _chunk_cache[cache_key] = text
        _chunk_cache.move_to_end(cache_key)
        while len(_chunk_cache) > API_CACHE_MAX_ENTRIES:
            _chunk_cache.popitem(last=False)

def append_generated_text(format_type, file_obj=None, memory_list=None, text=""):
    if not text:
        return
    if format_type in ['html', 'txt', 'md', 'epub']:
        text = sanitize_model_output(clean_text_artifacts(text), format_type)
        if not text:
            return
    if format_type in ['html', 'txt', 'md'] and file_obj:
        file_obj.write(text)
        file_obj.flush()
    if memory_list is not None:
        memory_list.append(text)

def batch_text_chunks(chunks, target_chars=TEXT_BATCH_TARGET_CHARS):
    batches = []
    current_batch = []
    current_len = 0
    for chunk in chunks:
        if not chunk:
            continue
        chunk_len = len(chunk)
        if current_batch and (current_len + chunk_len) > target_chars:
            batches.append("\n".join(current_batch))
            current_batch = [chunk]
            current_len = chunk_len
        else:
            current_batch.append(chunk)
            current_len += chunk_len
    if current_batch:
        batches.append("\n".join(current_batch))
    return batches

class HeartbeatMonitor:
    def __init__(self, timeout=HEARTBEAT_TIMEOUT_SEC):
        self.timeout = timeout
        self.timer = None
    def _stall_abort(self):
        print("\n[NETWORK STALL ALERT] API connection hung for over 5 minutes. Forcing exit to trigger fail-safe reboot.")
        import os
        os._exit(1)
    def ping(self):
        if self.timer: self.timer.cancel()
        self.timer = threading.Timer(self.timeout, self._stall_abort)
        self.timer.start()
    def stop(self):
        if self.timer: self.timer.cancel()

heartbeat = HeartbeatMonitor(HEARTBEAT_TIMEOUT_SEC)

def _filter_streaming_inline_image_payload(text, state):
    if not text:
        return ""
    if state is None:
        return text

    result = []
    idx = 0
    while idx < len(text):
        if state.get("suppressing"):
            quote = state.get("quote", '"')
            close_idx = text.find(quote, idx)
            if close_idx == -1:
                return "".join(result)
            result.append(quote)
            idx = close_idx + 1
            state["suppressing"] = False
            state["quote"] = '"'
            continue

        match = re.search(
            r'(?is)(src\s*=\s*)(["\'])(?:data:image\/[a-z0-9.+-]+;base64,|about:blank)',
            text[idx:],
        )
        if not match:
            result.append(text[idx:])
            break
        start = idx + match.start()
        end = idx + match.end()
        result.append(text[idx:start])
        result.append(f"{match.group(1)}{match.group(2)}about:blank")
        quote = match.group(2)
        close_idx = text.find(quote, end)
        if close_idx == -1:
            state["suppressing"] = True
            state["quote"] = quote
            return "".join(result)
        result.append(quote)
        idx = close_idx + 1
    return "".join(result)

def handle_stream(response, output_path, format_type, file_obj=None, memory_list=None, profile_key=None):
    global heartbeat
    heartbeat.ping()
    emitted = []
    total_chars = 0
    inline_image_state = {"suppressing": False, "quote": '"'} if format_type in ("html", "epub") else None
    try:
        for chunk in response:
            heartbeat.ping()
            text_val = ""
            if hasattr(chunk, 'text') and chunk.text: text_val = chunk.text
            elif hasattr(chunk, 'type') and chunk.type == 'content_block_delta': text_val = chunk.delta.text
            elif hasattr(chunk, 'choices') and len(chunk.choices) > 0 and chunk.choices[0].delta.content: text_val = chunk.choices[0].delta.content
            
            if text_val:
                if inline_image_state is not None:
                    text_val = _filter_streaming_inline_image_payload(text_val, inline_image_state)
                clean_chunk = sanitize_model_output(clean_text_artifacts(text_val), format_type)
                emitted.append(clean_chunk)
                total_chars += len(clean_chunk)
                append_generated_text(format_type, file_obj=file_obj, memory_list=memory_list, text=clean_chunk)
                if profile_key == "newspaper" and format_type == "html" and total_chars >= 12000:
                    notice = core_build_newspaper_safety_notice(format_type)
                    if notice:
                        append_generated_text(format_type, file_obj=file_obj, memory_list=memory_list, text="\n" + notice)
                        emitted.append("\n" + notice)
                    break
    finally:
        heartbeat.stop()
    return "".join(emitted)

def stream_with_cache(cache_key, request_fn, output_path, format_type, file_obj=None, memory_list=None, profile_key=None):
    cached = cache_get(cache_key)
    if cached is not None:
        append_generated_text(format_type, file_obj=file_obj, memory_list=memory_list, text=cached)
        return cached
    cleanup_fn = None
    request_result = request_fn()
    if (
        isinstance(request_result, tuple)
        and len(request_result) == 2
        and callable(request_result[1])
    ):
        response, cleanup_fn = request_result
    else:
        response = request_result
    try:
        generated = handle_stream(
            response,
            output_path,
            format_type,
            file_obj=file_obj,
            memory_list=memory_list,
            profile_key=profile_key,
        )
        cache_put(cache_key, generated)
        return generated
    finally:
        if cleanup_fn:
            cleanup_fn()

def play_completion_sound():
    print("\nProcessing complete!")
    if platform.system() == "Darwin":
        subprocess.Popen(["afplay", "/System/Library/Sounds/Glass.aiff"])
    elif platform.system() == "Windows":
        subprocess.Popen([
            "PowerShell",
            "-Command",
            "(New-Object Media.SoundPlayer 'C:\\Windows\\Media\\notify.wav').PlaySync();"
        ])

def process_files():
    config = get_user_preferences()
    model_name = config['model_name']
    api_key = get_api_key(model_name)
    try:
        if "claude" in model_name:
            import anthropic
            client = anthropic.Anthropic(api_key=api_key)
        elif "gpt" in model_name:
            import openai
            client = openai.OpenAI(api_key=api_key)
        else:
            client = genai.Client(api_key=api_key)
    except Exception as e:
        return print(f"Connection Error: {e}")

    prompt_text = get_prompt(config)
    target_scan_dir = BATCH_INPUT_DIR if config['batch_mode'].startswith('recursive') else INPUT_DIR
    os.makedirs(target_scan_dir, exist_ok=True)
    os.makedirs(config['output_dir'], exist_ok=True)
    valid_files = []
    if config['batch_mode'].startswith('recursive'):
        for r, d, f in os.walk(target_scan_dir):
            valid_files.extend([os.path.join(r, file) for file in f if os.path.splitext(file)[1].lower() in SUPPORTED_EXTENSIONS and not file.startswith('.')])
    else:
        if os.path.exists(target_scan_dir):
            valid_files = [os.path.join(target_scan_dir, f) for f in os.listdir(target_scan_dir) if os.path.splitext(f)[1].lower() in SUPPORTED_EXTENSIONS and not f.startswith('.')]
    if not valid_files:
        return print("No valid files found.")

    master_path = os.path.join(config['output_dir'], f"Chronicle_Merged.{config['format_type']}")
    master_file_obj, master_memory = None, []
    if config['merge_files']:
        if os.path.exists(master_path):
            os.remove(master_path)
        if config['format_type'] in ['html', 'txt', 'md']:
            master_file_obj = open(master_path, 'w', encoding='utf-8')
            write_header(master_file_obj, "Chronicle Merged", config['format_type'], get_output_lang_code(config), get_output_text_direction(config))

    for file_path in sorted(valid_files):
        filename = os.path.basename(file_path)
        base_name, ext = os.path.splitext(filename)[0], os.path.splitext(filename)[1].lower()
        current_memory = []
        if config['merge_files']:
            output_path, active_write_path = master_path, master_path
            current_file_obj, current_memory = master_file_obj, master_memory
            if config['format_type'] == 'html':
                current_file_obj.write("<br>")
            elif config['format_type'] in ['txt', 'md']:
                current_file_obj.write("\n\n")
        else:
            output_path = os.path.join(config['output_dir'], f"{base_name}.{config['format_type']}")
            if os.path.exists(output_path):
                if config.get('collision_mode') == 'skip':
                    continue
                elif config.get('collision_mode') == 'auto':
                    base_name = f"{base_name}_{int(time.time())}"
                    output_path = os.path.join(config['output_dir'], f"{base_name}.{config['format_type']}")
            active_write_path = output_path + ".tmp"
            if os.path.exists(active_write_path):
                os.remove(active_write_path)
            current_file_obj = open(active_write_path, 'w', encoding='utf-8') if config['format_type'] in ['html', 'txt', 'md'] else None
            if current_file_obj:
                write_header(
                    current_file_obj,
                    base_name,
                    config['format_type'],
                    get_output_lang_code(config),
                    get_output_text_direction(config),
                )
        print(f"\nProcessing {filename}...")
        try:
            config['source_path'] = file_path
            if ext == '.pdf':
                process_pdf(client, file_path, active_write_path, config['format_type'], prompt_text, config['model_name'], current_file_obj, current_memory)
            elif ext == '.pptx':
                process_pptx_document(client, file_path, active_write_path, config['format_type'], prompt_text, config['model_name'], current_file_obj, current_memory)
            elif ext in ['.docx', '.txt', '.md', '.rtf', '.csv', '.js', '.xlsx']:
                process_text_document(client, file_path, active_write_path, ext, config['format_type'], prompt_text, config['model_name'], current_file_obj, current_memory)
            elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.webp']:
                process_image(client, file_path, active_write_path, config['format_type'], prompt_text, config['model_name'], current_file_obj, current_memory)
            if config['batch_mode'] == 'recursive_delete':
                if is_path_within_protected_input_dirs(file_path):
                    print(f"Protected folder rule: not deleting source file in protected directory: {file_path}")
                else:
                    os.remove(file_path)
            if not config['merge_files']:
                dispatch_save(config, active_write_path, current_memory, base_name)
                if current_file_obj:
                    write_footer(current_file_obj, config['format_type'])
                    current_file_obj.close()
                if config['format_type'] == 'html' and os.path.exists(active_write_path):
                    try:
                        with open(active_write_path, 'r', encoding='utf-8', errors='ignore') as fh:
                            cleaned_html = normalize_streamed_html_document(fh.read())
                        if config.get('modernize_punctuation'):
                            cleaned_html = apply_modern_punctuation(cleaned_html)
                        if config.get('unit_conversion'):
                            cleaned_html = apply_modern_currency(cleaned_html)
                        if config.get('abbrev_expansion'):
                            cleaned_html = apply_expanded_abbreviations(cleaned_html)
                        cleaned_html = apply_newspaper_html_safety_fallback(
                            cleaned_html,
                            'html',
                            config.get('doc_profile'),
                        )
                        cleaned_html = recover_newspaper_header_citation(
                            cleaned_html,
                            'html',
                            config.get('doc_profile'),
                            config.get('source_path'),
                        )
                        cleaned_html = recover_source_attribution_footer(
                            cleaned_html,
                            'html',
                            config.get('doc_profile'),
                            config.get('source_path'),
                        )
                        cleaned_html = apply_handwriting_audit_flag(
                            cleaned_html,
                            'html',
                            config.get('doc_profile'),
                            whole_document=True,
                        )
                        cleaned_html = enforce_archival_heading_structure(
                            cleaned_html,
                            'html',
                            config.get('doc_profile'),
                        )
                        with open(active_write_path, 'w', encoding='utf-8') as fh:
                            fh.write(cleaned_html)
                    except Exception as cleanup_ex:
                        print(f"Warning: HTML cleanup skipped ({cleanup_ex})")
                elif config['format_type'] in ('txt', 'md') and os.path.exists(active_write_path):
                    try:
                        with open(active_write_path, 'r', encoding='utf-8', errors='ignore') as fh:
                            cleaned_text = fh.read()
                        if config.get('modernize_punctuation'):
                            cleaned_text = apply_modern_punctuation(cleaned_text)
                        if config.get('unit_conversion'):
                            cleaned_text = apply_modern_currency(cleaned_text)
                        if config.get('abbrev_expansion'):
                            cleaned_text = apply_expanded_abbreviations(cleaned_text)
                        cleaned_text = apply_handwriting_audit_flag(
                            cleaned_text,
                            config['format_type'],
                            config.get('doc_profile'),
                            whole_document=True,
                        )
                        with open(active_write_path, 'w', encoding='utf-8') as fh:
                            fh.write(cleaned_text)
                    except Exception as cleanup_ex:
                        print(f"Warning: text audit cleanup skipped ({cleanup_ex})")
                if os.path.exists(active_write_path):
                    os.rename(active_write_path, output_path)
            else:
                dispatch_save(config, master_path, master_memory, "Merged", clear_memory=True)
        except Exception as e:
            print(f"Error on {filename}: {e}")
            if current_file_obj and not config['merge_files']:
                current_file_obj.close()
    if config['merge_files']:
        dispatch_save(config, master_path, master_memory, "Merged", clear_memory=True)
        if master_file_obj:
            write_footer(master_file_obj, config['format_type'])
            master_file_obj.close()
            if config['format_type'] == 'html' and os.path.exists(master_path):
                try:
                    with open(master_path, 'r', encoding='utf-8', errors='ignore') as fh:
                        cleaned_html = normalize_streamed_html_document(fh.read())
                        if config.get('modernize_punctuation'):
                            cleaned_html = apply_modern_punctuation(cleaned_html)
                        if config.get('unit_conversion'):
                            cleaned_html = apply_modern_currency(cleaned_html)
                        if config.get('abbrev_expansion'):
                            cleaned_html = apply_expanded_abbreviations(cleaned_html)
                    cleaned_html = apply_newspaper_html_safety_fallback(
                        cleaned_html,
                        'html',
                        config.get('doc_profile'),
                    )
                    cleaned_html = recover_newspaper_header_citation(
                        cleaned_html,
                        'html',
                        config.get('doc_profile'),
                        config.get('source_path'),
                    )
                    cleaned_html = recover_source_attribution_footer(
                        cleaned_html,
                        'html',
                        config.get('doc_profile'),
                        config.get('source_path'),
                    )
                    cleaned_html = apply_handwriting_audit_flag(
                        cleaned_html,
                        'html',
                        config.get('doc_profile'),
                        whole_document=True,
                    )
                    cleaned_html = enforce_archival_heading_structure(
                        cleaned_html,
                        'html',
                        config.get('doc_profile'),
                    )
                    with open(master_path, 'w', encoding='utf-8') as fh:
                        fh.write(cleaned_html)
                except Exception as cleanup_ex:
                    print(f"Warning: merge HTML cleanup skipped ({cleanup_ex})")
            elif config['format_type'] in ('txt', 'md') and os.path.exists(master_path):
                try:
                    with open(master_path, 'r', encoding='utf-8', errors='ignore') as fh:
                        cleaned_text = fh.read()
                        if config.get('modernize_punctuation'):
                            cleaned_text = apply_modern_punctuation(cleaned_text)
                        if config.get('unit_conversion'):
                            cleaned_text = apply_modern_currency(cleaned_text)
                        if config.get('abbrev_expansion'):
                            cleaned_text = apply_expanded_abbreviations(cleaned_text)
                    cleaned_text = apply_handwriting_audit_flag(
                        cleaned_text,
                        config['format_type'],
                        config.get('doc_profile'),
                        whole_document=True,
                    )
                    with open(master_path, 'w', encoding='utf-8') as fh:
                        fh.write(cleaned_text)
                except Exception as cleanup_ex:
                    print(f"Warning: merge text audit cleanup skipped ({cleanup_ex})")
    play_completion_sound()

def get_pdf_chunk_pages(model_name, profile_key, total_pages, default_chunk_pages=PDF_CHUNK_PAGES, *, file_size_mb=None):
    model_name = str(model_name or "").lower()
    profile_key = str(profile_key or "")
    avg_page_mb = None
    if file_size_mb is not None and total_pages:
        try:
            avg_page_mb = float(file_size_mb) / max(1, int(total_pages))
        except (TypeError, ValueError):
            avg_page_mb = None
    if "claude" in model_name or "gpt" in model_name:
        return 1
    if profile_key == "newspaper" and total_pages >= 24:
        return 2
    if profile_key == "newspaper" and avg_page_mb is not None and avg_page_mb >= 0.9:
        return 1
    if profile_key in {"legal", "government"}:
        if total_pages >= 150:
            return 1
        if "gemini-2.5-pro" in model_name and total_pages >= 60:
            return 2
    if total_pages >= 60:
        return 2
    return min(2, max(1, int(default_chunk_pages)))


def process_pdf(client, pdf_path, output_path, format_type, prompt_text, model_name, file_obj, memory_list):
    profile_key = "newspaper" if "HISTORICAL NEWSPAPER RULES" in prompt_text else None
    if model_name.startswith("gemini-") and profile_key != "newspaper" and _probe_pdf_text_layer(pdf_path):
        legacy_process_pdf_gemini(client, pdf_path, format_type, prompt_text, model_name, file_obj, memory_list)
        return

    doc = fitz.open(pdf_path)
    try:
        total_pages = len(doc)
        current_page = 0

        while current_page < total_pages:
            image_bytes = _render_pdf_page_to_png_bytes(doc, current_page)
            temp_png = _write_temp_png("chronicle_cli_pdf_page_", image_bytes)
            page_fingerprint = sha256_text(image_bytes.hex())

            try:
                if "claude" in model_name or "gpt" in model_name:
                    payload = build_multimodal_payload(model_name, prompt_text, temp_png, "image/png")
                    cache_key = build_request_cache_key(
                        model_name, prompt_text, "pdf-vision", f"{current_page}:{page_fingerprint}"
                    )
                    stream_with_cache(
                        cache_key,
                        lambda: generate_with_retry(client, model_name, payload),
                        output_path,
                        format_type,
                        file_obj=file_obj,
                        memory_list=memory_list,
                        profile_key=profile_key,
                    )
                else:
                    cache_key = build_request_cache_key(
                        model_name, prompt_text, "pdf-upload", f"{current_page}:{page_fingerprint}"
                    )

                    def _gemini_pdf_request():
                        uploaded = client.files.upload(file=temp_png)
                        uploaded = wait_for_gemini_upload_ready(client, uploaded)
                        response = generate_with_retry(client, model_name, [uploaded, prompt_text])

                        def _cleanup_upload():
                            try:
                                client.files.delete(name=uploaded.name)
                            except Exception as cleanup_ex:
                                print(f"Warning: could not delete temporary Gemini upload {uploaded.name}: {cleanup_ex}")

                        return response, _cleanup_upload

                    stream_with_cache(
                        cache_key,
                        _gemini_pdf_request,
                        output_path,
                        format_type,
                        file_obj=file_obj,
                        memory_list=memory_list,
                        profile_key=profile_key,
                    )
            except Exception as e:
                print(f"\n[Gearshift Triggered] Error processing page {current_page + 1}: {e}")
                print(f"\n[CRITICAL FAILURE] Single page vision failed for page {current_page + 1}. Initiating sentence-level text fallback.")
                fallback_text = doc.load_page(current_page).get_text("text") or "[Unreadable Image Layer]"
                payload = fallback_text + "\n\n" + prompt_text if ("claude" in model_name or "gpt" in model_name) else [fallback_text, prompt_text]
                cache_key = build_request_cache_key(
                    model_name, prompt_text, "pdf-text-fallback", sha256_text(fallback_text)
                )
                stream_with_cache(
                    cache_key,
                    lambda: generate_with_retry(client, model_name, payload),
                    output_path,
                    format_type,
                    file_obj=file_obj,
                    memory_list=memory_list,
                    profile_key=profile_key,
                )
            finally:
                if os.path.exists(temp_png):
                    os.remove(temp_png)
            current_page += 1
            if current_page > 0 and current_page % 2 == 0:
                gc.collect()
    finally:
        doc.close()

def process_text_document(client, file_path, output_path, ext, format_type, prompt_text, model_name, file_obj, memory_list):
    profile_key = "newspaper" if "HISTORICAL NEWSPAPER RULES" in prompt_text else None
    is_tabular_html = "TABULAR DATA & SPREADSHEETS RULES" in prompt_text and format_type == "html"
    if is_tabular_html and ext in ('.csv', '.xlsx'):
        datasets = []
        if ext == '.csv':
            raw_csv = open(file_path, 'r', encoding='utf-8', errors='ignore').read()
            rows = core_parse_csv_rows(raw_csv)
            if rows:
                datasets.append({
                    "name": "Table Data",
                    "headers": [str(cell).strip() for cell in rows[0]],
                    "rows": [[str(cell).strip() for cell in row] for row in rows[1:]],
                })
        elif ext == '.xlsx' and openpyxl is not None:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for name in wb.sheetnames:
                rows = []
                for row in wb[name].iter_rows(values_only=True):
                    values = ["" if cell is None else str(cell).strip() for cell in row]
                    if any(values):
                        rows.append(values)
                if rows:
                    datasets.append({"name": name, "headers": rows[0], "rows": rows[1:]})
        if datasets:
            fragment = core_build_tabular_html_fragment(os.path.splitext(os.path.basename(file_path))[0], datasets)
            if file_obj:
                file_obj.write(fragment)
                file_obj.flush()
            if memory_list is not None:
                memory_list.append(fragment)
            return
    full_text = ""
    if ext == '.docx':
        full_text = "\n".join([p.text for p in docx.Document(file_path).paragraphs])
    elif ext == '.csv':
        raw_csv = open(file_path, 'r', encoding='utf-8', errors='ignore').read()
        full_text = csv_to_accessible_text(raw_csv)
    elif ext == '.xlsx':
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for name in wb.sheetnames:
            full_text += f"\n[--- Tab: {name} ---]\n"
            for row in wb[name].iter_rows(values_only=True):
                if any(c for c in row if str(c).strip()):
                    full_text += " | ".join([str(c) if c else "" for c in row]) + "\n"
    else:
        full_text = open(file_path, 'r', encoding='utf-8', errors='ignore').read()
    
    scrubbed_text = clean_text_artifacts(full_text)
    
    chunks, current = [], ""
    for p in scrubbed_text.split('\n'):
        if len(current) + len(p) > TEXT_CHUNK_CHARS and current:
            chunks.append(current)
            current = p + "\n"
        else:
            current += p + "\n"
    if current:
        chunks.append(current)
    batches = batch_text_chunks(chunks)
    for chunk in batches:
        payload = chunk + "\n\n" + prompt_text if ("claude" in model_name or "gpt" in model_name) else [chunk, prompt_text]
        cache_key = build_request_cache_key(model_name, prompt_text, "text", sha256_text(chunk))
        stream_with_cache(
            cache_key,
            lambda: generate_with_retry(client, model_name, payload),
            output_path,
            format_type,
            file_obj=file_obj,
            memory_list=memory_list,
            profile_key=profile_key,
        )

def process_image(client, file_path, output_path, format_type, prompt_text, model_name, file_obj, memory_list):
    profile_key = "newspaper" if "HISTORICAL NEWSPAPER RULES" in prompt_text else None
    enhanced_path = enhance_image_for_microtext(file_path)
    if "claude" in model_name or "gpt" in model_name:
        payload = build_multimodal_payload(model_name, prompt_text, enhanced_path, "image/png")
        cache_key = build_request_cache_key(model_name, prompt_text, "image", sha256_file(enhanced_path))
        stream_with_cache(
            cache_key,
            lambda: generate_with_retry(client, model_name, payload),
            output_path,
            format_type,
            file_obj=file_obj,
            memory_list=memory_list,
            profile_key=profile_key,
        )
    else:
        cache_key = build_request_cache_key(model_name, prompt_text, "image-upload", sha256_file(enhanced_path))
        def _gemini_image_request():
            uploaded = client.files.upload(file=enhanced_path)
            uploaded = wait_for_gemini_upload_ready(client, uploaded, poll_sec=2.0)
            response = generate_with_retry(client, model_name, [uploaded, prompt_text])
            def _cleanup_upload():
                try:
                    client.files.delete(name=uploaded.name)
                except Exception as cleanup_ex:
                    print(f"Warning: could not delete temporary Gemini upload {uploaded.name}: {cleanup_ex}")
            return response, _cleanup_upload
        stream_with_cache(
            cache_key,
            _gemini_image_request,
            output_path,
            format_type,
            file_obj=file_obj,
            memory_list=memory_list,
            profile_key=profile_key,
        )
    if enhanced_path != file_path:
        os.remove(enhanced_path)

def wait_for_gemini_upload_ready(client, uploaded, poll_sec=GEMINI_UPLOAD_POLL_SEC, max_wait_sec=GEMINI_UPLOAD_MAX_WAIT_SEC, time_fn=time.time, sleep_fn=time.sleep, log_cb=print):
    start = time_fn()
    next_log_at = 0.0
    current = uploaded
    while getattr(getattr(current, "state", None), "name", "") == "PROCESSING":
        elapsed = max(0.0, time_fn() - start)
        if elapsed >= max_wait_sec:
            raise TimeoutError(
                f"Google file preparation exceeded {int(max_wait_sec)} seconds for {getattr(current, 'name', 'upload')}."
            )
        if elapsed >= next_log_at:
            log_cb(
                f"[Upload Processing] Waiting for Google to prepare {getattr(current, 'name', 'upload')} ({elapsed:.0f}s elapsed)..."
            )
            next_log_at += 15.0
        sleep_fn(min(poll_sec, max_wait_sec - elapsed))
        current = client.files.get(name=current.name)
    if getattr(getattr(current, "state", None), "name", "") == "FAILED":
        raise Exception("Google API failed to process the uploaded document.")
    return current

def generate_with_retry(client, model_name, contents, max_retries=5, base_delay=10):
    import time
    for attempt in range(max_retries):
        try:
            wait_for_request_slot()
            try:
                pace_api_request()
                if "claude" in model_name:
                    return client.messages.create(
                        model=model_name,
                        max_tokens=8192,
                        messages=[{"role": "user", "content": contents}],
                        stream=True,
                        extra_headers={"anthropic-beta": "pdfs-2024-09-25"},
                    )
                if "gpt" in model_name:
                    return client.chat.completions.create(
                        model=model_name,
                        messages=[{"role": "user", "content": contents}],
                        stream=True,
                    )
                return client.models.generate_content_stream(model=model_name, contents=contents)
            finally:
                release_request_slot()
        except Exception as e:
            error_str = str(e).lower()
            if "401" in error_str or "unauthorized" in error_str or "invalid api key" in error_str or "authentication" in error_str:
                raise Exception("Authentication failed (401/invalid API key). Check your API key for the selected engine.")
            if is_connection_path_error(e):
                print(
                    f"\n[Polite Backoff] Connection path unavailable. Waiting {API_CONNECTION_FAILURE_DELAY_SEC:.0f}s "
                    "before giving up to avoid hammering the service."
                )
                time.sleep(API_CONNECTION_FAILURE_DELAY_SEC)
                raise Exception(
                    "Connection path unavailable. Chronicle stopped retrying to avoid hammering the API service."
                ) from e
            if "429" in error_str or "exhausted" in error_str or "quota" in error_str or "overloaded" in error_str:
                delay = base_delay * (2 ** attempt)
                print(f"\n[API Rate Limit Hit] Backing off for {delay} seconds (Attempt {attempt + 1}/{max_retries})...")
                time.sleep(delay)
            else:
                raise e
    raise Exception("Max retries exceeded. API is fundamentally unresponsive.")

def build_multimodal_payload(model_name, prompt_text, file_path=None, mime_type="image/png"):
    import base64
    if not file_path:
        return prompt_text
    with open(file_path, "rb") as f:
        b64_data = base64.b64encode(f.read()).decode('utf-8')
    if "claude" in model_name:
        if "pdf" in mime_type:
            return [
                {"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": b64_data}},
                {"type": "text", "text": prompt_text},
            ]
        return [
            {"type": "image", "source": {"type": "base64", "media_type": mime_type, "data": b64_data}},
            {"type": "text", "text": prompt_text},
        ]
    if "gpt" in model_name:
        if "pdf" in mime_type:
            raise Exception("GPT-4o engine rejected Base64 PDF. Forcing sentence-level text fallback.")
        return [
            {"type": "text", "text": prompt_text},
            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{b64_data}"}},
        ]
    return None

def process_pptx_document(client, file_path, output_path, format_type, prompt_text, model_name, file_obj, memory_list):
    import pptx
    try:
        from pptx.enum.shape import MSO_SHAPE_TYPE
    except ImportError:
        MSO_SHAPE_TYPE = None

    # Gemini can process native PPTX uploads; use this first to preserve visual context.
    if "claude" not in model_name and "gpt" not in model_name:
        upload_cache_key = build_request_cache_key(model_name, prompt_text, "pptx-upload", sha256_file(file_path))
        try:
            def _gemini_pptx_request():
                uploaded = client.files.upload(file=file_path)
                try:
                    while getattr(getattr(uploaded, "state", None), "name", "") == "PROCESSING":
                        time.sleep(2)
                        uploaded = client.files.get(name=uploaded.name)
                    if getattr(getattr(uploaded, "state", None), "name", "") == "FAILED":
                        raise Exception("Google API failed to process the PPTX document.")
                except Exception as poll_ex:
                    print(f"Warning: PPTX upload polling failed ({poll_ex}); attempting generation anyway.")
                response = generate_with_retry(client, model_name, [uploaded, prompt_text])
                def _cleanup_upload():
                    try:
                        client.files.delete(name=uploaded.name)
                    except Exception as cleanup_ex:
                        print(f"Warning: could not delete temporary Gemini upload {uploaded.name}: {cleanup_ex}")
                return response, _cleanup_upload

            stream_with_cache(
                upload_cache_key,
                _gemini_pptx_request,
                output_path,
                format_type,
                file_obj=file_obj,
                memory_list=memory_list,
            )
            return
        except Exception as upload_ex:
            print(f"Warning: native PPTX vision failed ({upload_ex}); falling back to structured slide extraction.")

    prs = pptx.Presentation(file_path)
    full_text = ""
    slide_num = 1
    for slide in prs.slides:
        full_text += f"\n[--- Slide: {slide_num} ---]\n"
        try:
            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes_text:
                    full_text += f"[Slide Notes]\n{notes_text}\n"
        except Exception:
            pass
        try:
            shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left))
        except Exception:
            shapes = slide.shapes
        for shape_idx, shape in enumerate(shapes, start=1):
            shape_name = (getattr(shape, "name", "") or "").strip()
            shape_label = f"Slide {slide_num} Shape {shape_idx}"
            if shape_name:
                shape_label += f" ({shape_name})"

            if getattr(shape, "has_text_frame", False):
                shape_text = (getattr(shape, "text", "") or "").strip()
                if shape_text:
                    full_text += shape_text + "\n"

            if getattr(shape, "has_table", False):
                full_text += f"[Table Object: {shape_label}]\n"
                try:
                    for row in shape.table.rows:
                        row_vals = []
                        for cell in row.cells:
                            row_vals.append((cell.text or "").replace("\n", " ").strip())
                        full_text += " | ".join(row_vals) + "\n"
                except Exception:
                    full_text += "[Table content unreadable]\n"

            if getattr(shape, "has_chart", False):
                full_text += f"[Chart Object: {shape_label}]\n"
                try:
                    chart = shape.chart
                    full_text += f"Chart Type: {chart.chart_type}\n"
                    try:
                        if chart.has_title and chart.chart_title and chart.chart_title.text_frame:
                            chart_title = (chart.chart_title.text_frame.text or "").strip()
                            if chart_title:
                                full_text += f"Chart Title: {chart_title}\n"
                    except Exception:
                        pass
                    categories = []
                    try:
                        categories = [str(cat) for cat in chart.plots[0].categories]
                    except Exception:
                        categories = []
                    if categories:
                        full_text += "Categories: " + " | ".join(categories) + "\n"
                    for series_idx, series in enumerate(chart.series, start=1):
                        try:
                            series_name = str(series.name)
                        except Exception:
                            series_name = f"Series {series_idx}"
                        try:
                            values = [str(v) for v in series.values]
                        except Exception:
                            values = []
                        if values:
                            full_text += f"{series_name}: " + " | ".join(values) + "\n"
                        else:
                            full_text += f"{series_name}: [No numeric values extracted]\n"
                except Exception:
                    full_text += "[Chart metadata unreadable]\n"

            shape_type = getattr(shape, "shape_type", None)
            visual_shape_types = {
                getattr(MSO_SHAPE_TYPE, "GROUP", None),
                getattr(MSO_SHAPE_TYPE, "DIAGRAM", None),
                getattr(MSO_SHAPE_TYPE, "CANVAS", None),
                getattr(MSO_SHAPE_TYPE, "LINKED_PICTURE", None),
            }
            if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.PICTURE:
                full_text += f"[Image Object: {shape_label}]\n"
            elif shape_type in visual_shape_types:
                full_text += f"[Visual Object: {shape_label}; Type={shape_type}]\n"
        slide_num += 1
    
    scrubbed_input = clean_text_artifacts(full_text)
    chunks = [scrubbed_input[i:i+TEXT_CHUNK_CHARS] for i in range(0, len(scrubbed_input), TEXT_CHUNK_CHARS)]
    
    batches = batch_text_chunks(chunks)
    for chunk in batches:
        payload = chunk + "\n\n" + prompt_text if ("claude" in model_name or "gpt" in model_name) else [chunk, prompt_text]
        cache_key = build_request_cache_key(model_name, prompt_text, "pptx", sha256_text(chunk))
        stream_with_cache(
            cache_key,
            lambda: generate_with_retry(client, model_name, payload),
            output_path,
            format_type,
            file_obj=file_obj,
            memory_list=memory_list,
        )

if __name__ == "__main__":
    migrate_legacy_cli_key_files()
    shared_emit_launch_continuity(script_dir=SCRIPT_DIR)
    process_files()
