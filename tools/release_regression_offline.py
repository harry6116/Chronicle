#!/usr/bin/env python3
from __future__ import annotations

import json
import os
import re
import shutil
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from types import SimpleNamespace

import docx
import openpyxl
from PIL import Image, ImageDraw
from pypdf import PdfWriter
from ebooklib import epub
import pptx

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

import chronicle_runtime as chronicle


@dataclass
class Case:
    input_path: Path
    output_format: str
    name: str


class FakeChunk:
    def __init__(self, text: str):
        self.text = text


class FakeFilesApi:
    def upload(self, file: str):
        name = f"fake_{Path(file).name}"
        return SimpleNamespace(name=name, state=SimpleNamespace(name="DONE"))

    def get(self, name: str):
        return SimpleNamespace(name=name, state=SimpleNamespace(name="DONE"))

    def delete(self, name: str):
        return None


class FakeClient:
    def __init__(self):
        self.files = FakeFilesApi()


def fake_generate_with_retry(client, model_name, contents, max_retries=5, base_delay=10):
    txt = (
        "# Synthetic Heading\n"
        "Recovered sample paragraph.\n"
        "data:image/png;base64,AAAA\n"
        "<table><tr><td>x</td></tr></table>\n"
        "```html\n<figure>placeholder</figure>\n```\n"
        "Row 1: a=1 | b=2\n"
    )
    return [FakeChunk(txt)]


def write_pdf(path: Path):
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with open(path, "wb") as f:
        writer.write(f)


def write_docx(path: Path):
    d = docx.Document()
    d.add_heading("Sample DOCX", level=1)
    d.add_paragraph("Body paragraph for regression testing.")
    d.save(path)


def write_txt(path: Path):
    path.write_text("Plain text line one.\nLine two.", encoding="utf-8")


def write_md(path: Path):
    path.write_text("# Markdown Input\n\n- item one\n- item two", encoding="utf-8")


def write_rtf(path: Path):
    path.write_text(r"{\rtf1\ansi Sample RTF content for test.}", encoding="utf-8")


def write_csv(path: Path):
    path.write_text("col_a,col_b\n1,2\n3,4\n", encoding="utf-8")


def write_xlsx(path: Path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["name", "value"])
    ws.append(["alpha", 10])
    ws.append(["bravo", 20])
    wb.save(path)


def write_pptx(path: Path):
    pres = pptx.Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Bullet one\nBullet two"
    pres.save(path)


def write_image(path: Path):
    img = Image.new("RGB", (640, 360), color=(245, 245, 245))
    dr = ImageDraw.Draw(img)
    dr.text((20, 20), "Image input sample", fill=(0, 0, 0))
    img.save(path)


def write_epub(path: Path):
    book = epub.EpubBook()
    book.set_identifier("offline-regression-epub")
    book.set_title("Regression EPUB")
    book.set_language("en")
    ch = epub.EpubHtml(title="Chapter 1", file_name="chap1.xhtml", lang="en")
    ch.content = "<h1>Chapter 1</h1><p>EPUB input content.</p>"
    book.add_item(ch)
    book.spine = ["nav", ch]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    epub.write_epub(str(path), book, {})


def make_inputs(inp_dir: Path) -> list[Path]:
    files = []
    creators = [
        ("sample_01.pdf", write_pdf),
        ("sample_02.docx", write_docx),
        ("sample_03.txt", write_txt),
        ("sample_04.md", write_md),
        ("sample_05.rtf", write_rtf),
        ("sample_06.csv", write_csv),
        ("sample_07.xlsx", write_xlsx),
        ("sample_08.pptx", write_pptx),
        ("sample_09.jpg", write_image),
        ("sample_10.epub", write_epub),
    ]
    for name, fn in creators:
        p = inp_dir / name
        fn(p)
        files.append(p)
    return files


def run_case(case: Case, out_dir: Path, client: FakeClient) -> tuple[bool, str]:
    cfg = {
        "format_type": case.output_format,
        "model_name": "gemini-2.5-pro",
        "doc_profile": "standard",
        "translate_mode": "none",
        "translate_target": "English",
        "modernize_punctuation": False,
        "unit_conversion": False,
        "merge_files": False,
        "collision_mode": "overwrite",
        "image_descriptions": True,
        "large_print": False,
        "batch_mode": "flat",
        "abbrev_expansion": False,
        "academic_mode": False,
    }
    prompt = chronicle.get_prompt(cfg)
    base = case.input_path.stem
    out_path = out_dir / f"{base}.{case.output_format}"
    tmp_out = Path(str(out_path) + ".tmp")
    if tmp_out.exists():
        tmp_out.unlink()
    memory = []
    file_obj = None

    try:
        if case.output_format in {"html", "txt", "md"}:
            file_obj = open(tmp_out, "w", encoding="utf-8")
            chronicle.write_header(
                file_obj,
                base,
                case.output_format,
                chronicle.get_output_lang_code(cfg),
                chronicle.get_output_text_direction(cfg),
            )
        ext = case.input_path.suffix.lower()
        if ext == ".pdf":
            chronicle.process_pdf(client, str(case.input_path), str(tmp_out), case.output_format, prompt, cfg["model_name"], file_obj, memory)
        elif ext == ".pptx":
            chronicle.process_pptx_document(client, str(case.input_path), str(tmp_out), case.output_format, prompt, cfg["model_name"], file_obj, memory)
        elif ext in {".docx", ".txt", ".md", ".rtf", ".csv", ".js", ".xlsx"}:
            chronicle.process_text_document(client, str(case.input_path), str(tmp_out), ext, case.output_format, prompt, cfg["model_name"], file_obj, memory)
        elif ext in {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".webp"}:
            chronicle.process_image(client, str(case.input_path), str(tmp_out), case.output_format, prompt, cfg["model_name"], file_obj, memory)
        elif ext == ".epub":
            chronicle.process_text_document(client, str(case.input_path), str(tmp_out), ext, case.output_format, prompt, cfg["model_name"], file_obj, memory)
        else:
            return False, f"unsupported test extension {ext}"

        chronicle.dispatch_save(cfg, str(tmp_out), memory, base)
        if file_obj:
            chronicle.write_footer(file_obj, case.output_format)
            file_obj.close()
        if case.output_format == "html" and tmp_out.exists():
            cleaned = chronicle.normalize_streamed_html_document(tmp_out.read_text(encoding="utf-8", errors="ignore"))
            cleaned = chronicle.enforce_archival_heading_structure(cleaned, "html", cfg.get("doc_profile"))
            tmp_out.write_text(cleaned, encoding="utf-8")
        if tmp_out.exists():
            os.replace(tmp_out, out_path)

        if not out_path.exists() or out_path.stat().st_size == 0:
            return False, "missing or empty output"

        # Minimal sanitizer assertions for non-HTML text-like outputs.
        if case.output_format in {"txt", "md", "csv", "json"}:
            payload = out_path.read_text(encoding="utf-8", errors="ignore")
            if "data:image/" in payload.lower():
                return False, "base64 leak detected"
            if "```" in payload:
                return False, "fence wrapper leak detected"
            if re.search(r"</?(?:html|head|body|table|tr|td|th|div|span|figure|figcaption)\b", payload, flags=re.IGNORECASE):
                return False, "html tag leak detected"

        return True, str(out_path)
    except Exception as ex:
        return False, str(ex)
    finally:
        if file_obj and not file_obj.closed:
            file_obj.close()
        if tmp_out.exists():
            try:
                tmp_out.unlink()
            except Exception:
                pass


def main() -> int:
    chronicle.generate_with_retry = fake_generate_with_retry  # type: ignore
    with tempfile.TemporaryDirectory(prefix="chronicle_release_regression_") as td:
        root = Path(td)
        inp_dir = root / "inputs"
        out_dir = root / "outputs"
        inp_dir.mkdir(parents=True, exist_ok=True)
        out_dir.mkdir(parents=True, exist_ok=True)
        inputs = make_inputs(inp_dir)
        formats = ["html", "txt", "docx", "md", "pdf", "json", "csv", "epub", "txt", "html"]
        cases = [Case(path, fmt, f"{path.name}->{fmt}") for path, fmt in zip(inputs, formats)]

        client = FakeClient()
        report = {"total": len(cases), "passed": 0, "failed": 0, "results": []}
        for c in cases:
            ok, detail = run_case(c, out_dir, client)
            report["results"].append({"case": c.name, "ok": ok, "detail": detail})
            if ok:
                report["passed"] += 1
            else:
                report["failed"] += 1

        print(json.dumps(report, indent=2))
        return 0 if report["failed"] == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
