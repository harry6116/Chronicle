import os
import re
import subprocess
import tempfile
from email import policy
from email.parser import BytesParser

from chronicle_core import build_tabular_html_fragment, parse_csv_rows
from chronicle_app.services.runtime_policies import wait_for_gemini_upload_ready

try:  # pragma: no cover - optional import is environment dependent
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None

try:  # pragma: no cover - optional import is environment dependent
    from pillow_heif import register_heif_opener
except ImportError:  # pragma: no cover
    register_heif_opener = None


GEMINI_FILE_UPLOAD_TIMEOUT_MS = 300_000
GEMINI_FILE_DELETE_TIMEOUT_MS = 30_000
HEIC_IMAGE_EXTENSIONS = {".heic", ".heif"}
PIL_STAGED_IMAGE_EXTENSIONS = {".gif", ".avif", ".jp2", ".j2k", ".ppm", ".pgm", ".pbm"}
STANDARD_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".webp"}
SUPPORTED_IMAGE_EXTENSIONS = STANDARD_IMAGE_EXTENSIONS | HEIC_IMAGE_EXTENSIONS
SUPPORTED_IMAGE_EXTENSIONS |= PIL_STAGED_IMAGE_EXTENSIONS
TABULAR_TEXT_EXTENSIONS = {".csv", ".tsv", ".xlsx", ".xls"}
SOURCE_TEXT_EXTENSIONS = {
    ".docx", ".txt", ".md", ".rtf", ".csv", ".tsv", ".js", ".json", ".jsonl", ".xml",
    ".log", ".xlsx", ".xls", ".html", ".htm", ".eml",
}
EPUB_EXTENSIONS = {".epub"}
FITZ_TEXT_EXTENSIONS = {".mobi", ".fb2"}
PRESENTATION_EXTENSIONS = {".pptx", ".ppt"}
RENDERABLE_DOCUMENT_EXTENSIONS = {".svg", ".xps", ".oxps", ".cbz"}
SUPPORTED_EXTENSIONS = (
    {".pdf"}
    | SOURCE_TEXT_EXTENSIONS
    | EPUB_EXTENSIONS
    | FITZ_TEXT_EXTENSIONS
    | PRESENTATION_EXTENSIONS
    | SUPPORTED_IMAGE_EXTENSIONS
    | RENDERABLE_DOCUMENT_EXTENSIONS
)
RESUMABLE_UNIT_EXTENSIONS = (
    SOURCE_TEXT_EXTENSIONS
    | EPUB_EXTENSIONS
    | FITZ_TEXT_EXTENSIONS
    | PRESENTATION_EXTENSIONS
    | RENDERABLE_DOCUMENT_EXTENSIONS
)


def register_optional_image_decoders():
    if register_heif_opener is not None:
        register_heif_opener()


register_optional_image_decoders()


def supported_extensions_list(*, include_xls=True, include_legacy_ppt=True):
    excluded = set()
    if not include_xls:
        excluded.add(".xls")
    if not include_legacy_ppt:
        excluded.add(".ppt")
    return sorted(SUPPORTED_EXTENSIONS - excluded)


def supported_files_wildcard(*, include_xls=True, include_legacy_ppt=True):
    extensions = supported_extensions_list(include_xls=include_xls, include_legacy_ppt=include_legacy_ppt)
    return "Supported Files|" + ";".join(f"*{ext}" for ext in extensions)


def _make_temp_png(*, mkstemp_fn=tempfile.mkstemp, close_fn=os.close):
    fd, converted_path = mkstemp_fn(prefix="chronicle_image_", suffix=".png")
    close_fn(fd)
    return converted_path


def convert_heic_to_png_for_scan(
    path,
    *,
    log_cb=None,
    subprocess_module=subprocess,
    mkstemp_fn=tempfile.mkstemp,
    close_fn=os.close,
    remove_fn=os.remove,
):
    ext = os.path.splitext(path)[1].lower()
    if ext not in HEIC_IMAGE_EXTENSIONS:
        return path

    if Image is not None and register_heif_opener is not None:
        converted_path = _make_temp_png(mkstemp_fn=mkstemp_fn, close_fn=close_fn)
        try:
            with Image.open(path) as img:
                if getattr(img, "n_frames", 1) > 1:
                    img.seek(0)
                if img.mode not in ("RGB", "L"):
                    img = img.convert("RGB")
                img.save(converted_path, "PNG")
            if log_cb:
                log_cb(f"[Image] Converted {os.path.basename(path)} from HEIC/HEIF to PNG for scanning.")
            return converted_path
        except Exception:
            try:
                remove_fn(converted_path)
            except Exception:
                pass

    fd, converted_path = mkstemp_fn(prefix="chronicle_heic_", suffix=".png")
    close_fn(fd)
    try:
        proc = subprocess_module.run(
            ["sips", "-s", "format", "png", path, "--out", converted_path],
            capture_output=True,
            text=True,
            check=False,
        )
    except FileNotFoundError as ex:
        try:
            remove_fn(converted_path)
        except Exception:
            pass
        raise RuntimeError(
            "HEIC/HEIF image support requires macOS image conversion support. "
            "Convert this file to PNG or JPEG first on this system."
        ) from ex

    if proc.returncode != 0 or not os.path.exists(converted_path) or os.path.getsize(converted_path) == 0:
        try:
            remove_fn(converted_path)
        except Exception:
            pass
        detail = (proc.stderr or proc.stdout or "unknown conversion error").strip()
        raise RuntimeError(f"Could not convert HEIC/HEIF image for scanning: {detail}")

    if log_cb:
        log_cb(f"[Image] Converted {os.path.basename(path)} from HEIC/HEIF to PNG for scanning.")
    return converted_path


def convert_pil_image_to_png_for_scan(
    path,
    *,
    log_cb=None,
    image_module=Image,
    mkstemp_fn=tempfile.mkstemp,
    close_fn=os.close,
    remove_fn=os.remove,
):
    ext = os.path.splitext(path)[1].lower()
    if ext not in PIL_STAGED_IMAGE_EXTENSIONS:
        return path
    if image_module is None:
        raise RuntimeError(f"{ext.upper()} image support requires Pillow in the bundled runtime.")

    converted_path = _make_temp_png(mkstemp_fn=mkstemp_fn, close_fn=close_fn)
    try:
        with image_module.open(path) as img:
            if getattr(img, "n_frames", 1) > 1:
                img.seek(0)
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")
            img.save(converted_path, "PNG")
    except Exception as ex:
        try:
            remove_fn(converted_path)
        except Exception:
            pass
        raise RuntimeError(f"Could not convert {ext.upper()} image for scanning: {ex}") from ex

    if log_cb:
        log_cb(f"[Image] Converted {os.path.basename(path)} from {ext.upper()} to PNG for scanning.")
    return converted_path


def prepare_image_for_scan(
    path,
    *,
    log_cb=None,
    remove_fn=os.remove,
    image_module=Image,
    mkstemp_fn=tempfile.mkstemp,
    close_fn=os.close,
):
    heic_path = convert_heic_to_png_for_scan(
        path,
        log_cb=log_cb,
        remove_fn=remove_fn,
        mkstemp_fn=mkstemp_fn,
        close_fn=close_fn,
    )
    try:
        staged_path = convert_pil_image_to_png_for_scan(
            heic_path,
            log_cb=log_cb,
            remove_fn=remove_fn,
            image_module=image_module,
            mkstemp_fn=mkstemp_fn,
            close_fn=close_fn,
        )
    except Exception:
        if heic_path != path:
            try:
                remove_fn(heic_path)
            except Exception:
                pass
        raise
    if heic_path != path and staged_path != heic_path:
        try:
            remove_fn(heic_path)
        except Exception:
            pass
    return staged_path


def _strip_html_text(raw):
    try:
        from bs4 import BeautifulSoup  # type: ignore

        soup = BeautifulSoup(raw, "html.parser")
        for node in soup(["script", "style", "noscript"]):
            node.decompose()
        return soup.get_text("\n")
    except Exception:
        stripped = re.sub(r"(?is)<(script|style|noscript)\b[^>]*>.*?</\1>", " ", raw)
        return re.sub(r"<[^>]+>", " ", stripped)


def _load_email_content(path):
    with open(path, "rb") as fh:
        msg = BytesParser(policy=policy.default).parse(fh)
    parts = [
        f"Subject: {msg.get('subject', '')}",
        f"From: {msg.get('from', '')}",
        f"To: {msg.get('to', '')}",
        f"Date: {msg.get('date', '')}",
        "",
    ]
    body = ""
    if msg.is_multipart():
        html_fallback = ""
        for part in msg.walk():
            content_type = part.get_content_type()
            if part.get_content_disposition() == "attachment":
                continue
            if content_type == "text/plain" and not body:
                body = part.get_content()
            elif content_type == "text/html" and not html_fallback:
                html_fallback = _strip_html_text(part.get_content())
        if not body:
            body = html_fallback
    else:
        payload = msg.get_content()
        body = _strip_html_text(payload) if msg.get_content_type() == "text/html" else payload
    parts.append(body or "")
    return "\n".join(parts)


def _load_fitz_text_content(path, *, fitz_module):
    if fitz_module is None:
        import fitz as fitz_module  # type: ignore

    doc = fitz_module.open(path)
    try:
        full = ""
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text") if hasattr(page, "get_text") else ""
            if text.strip():
                full += f"\n[--- Page: {page_num + 1} ---]\n{text}\n"
        return full
    finally:
        doc.close()


def _load_text_document_content(
    path,
    ext,
    *,
    csv_to_accessible_text_fn,
    docx_module,
    openpyxl_module,
    subprocess_module=subprocess,
):
    full = ""
    if ext == ".docx":
        full = "\n".join([p.text for p in docx_module.Document(path).paragraphs])
    elif ext in {".csv", ".tsv"}:
        with open(path, "r", encoding="utf-8", errors="ignore") as fh:
            raw_csv = fh.read()
        delimiter = "\t" if ext == ".tsv" else ","
        full = csv_to_accessible_text_fn(raw_csv.replace("\t", ",") if delimiter == "\t" else raw_csv)
    elif ext in [".xlsx", ".xls"]:
        if ext == ".xls":
            try:
                import xlrd  # type: ignore

                wb_xls = xlrd.open_workbook(path)
                for sheet in wb_xls.sheets():
                    full += f"\n[--- Tab: {sheet.name} ---]\n"
                    for ridx in range(sheet.nrows):
                        row = sheet.row_values(ridx)
                        if any(str(c).strip() for c in row):
                            full += " | ".join([str(c) if c is not None else "" for c in row]) + "\n"
            except Exception:
                try:
                    conv = subprocess_module.run(
                        ["textutil", "-convert", "txt", "-stdout", path],
                        capture_output=True,
                        text=True,
                        check=False,
                    )
                    if conv.returncode == 0 and conv.stdout.strip():
                        full = conv.stdout
                    else:
                        raise RuntimeError("textutil conversion returned no content.")
                except Exception as ex:
                    raise Exception(
                        "Legacy .xls parsing failed. Install xlrd for robust .xls support or convert to .xlsx. "
                        f"Details: {ex}"
                    )
        else:
            workbook = openpyxl_module.load_workbook(path, data_only=True)
            for name in workbook.sheetnames:
                full += f"\n[--- Tab: {name} ---]\n"
                for row in workbook[name].iter_rows(values_only=True):
                    if any(c for c in row if str(c).strip()):
                        full += " | ".join([str(c) if c else "" for c in row]) + "\n"
    elif ext in {".html", ".htm"}:
        with open(path, "r", encoding="utf-8", errors="ignore") as fh:
            full = _strip_html_text(fh.read())
    elif ext == ".eml":
        full = _load_email_content(path)
    else:
        with open(path, "r", encoding="utf-8", errors="ignore") as fh:
            full = fh.read()
    return full


def _write_direct_output(text, *, f_obj, mem):
    if f_obj is not None:
        f_obj.write(text)
        f_obj.flush()
    if mem is not None:
        mem.append(text)


def _build_source_text_prompt(prompt, *, source_kind, doc_profile=None):
    profile = str(doc_profile or "standard").lower()
    mode = (
        "\n\nSOURCE-TEXT REPLICATION MODE:\n"
        f"- This input comes from Chronicle's direct {source_kind} text extraction path, not a visual scan.\n"
        "- Treat the supplied source text as authoritative unless it is clearly malformed or structurally broken.\n"
        "- Prioritize faithful replication of the original wording, numbering, hierarchy, and reading order.\n"
        "- Preserve clause numbering, headings, list nesting, table relationships, and explicit page or slide references when present.\n"
        "- Do not add visual-guess OCR repairs that are only appropriate for faded, skewed, or image-only scans.\n"
        "- Apply accessibility cleanup conservatively: fix only genuine structural/accessibility defects, not normal born-digital wording or punctuation.\n"
    )
    if profile in {"legal", "government"}:
        mode += (
            "- Be especially strict about legal/government hierarchy, cross-references, schedules, tables, and citation fidelity.\n"
        )
    return prompt + mode


def _build_tabular_datasets(path, ext, *, openpyxl_module):
    if ext in {".csv", ".tsv"}:
        with open(path, "r", encoding="utf-8", errors="ignore") as fh:
            raw = fh.read()
            rows = parse_csv_rows(raw.replace("\t", ",") if ext == ".tsv" else raw)
        if not rows:
            return []
        headers = [str(cell).strip() for cell in rows[0]]
        body = [[str(cell).strip() for cell in row] for row in rows[1:]]
        return [{"name": "Table Data", "headers": headers, "rows": body}]

    if ext in [".xlsx", ".xls"] and openpyxl_module is not None:
        workbook = openpyxl_module.load_workbook(path, data_only=True)
        datasets = []
        for name in workbook.sheetnames:
            rows = []
            for row in workbook[name].iter_rows(values_only=True):
                values = ["" if cell is None else str(cell).strip() for cell in row]
                if any(values):
                    rows.append(values)
            if not rows:
                continue
            datasets.append({
                "name": name,
                "headers": rows[0],
                "rows": rows[1:],
            })
        return datasets
    return []


def _stream_text_batches(
    full_text,
    *,
    out,
    fmt,
    prompt,
    model,
    f_obj,
    mem,
    log_cb,
    pause_cb,
    payload_kind,
    text_chunk_chars,
    clean_text_fn,
    batch_text_chunks_fn,
    build_request_cache_key_fn,
    sha256_text_fn,
    stream_with_cache_fn,
    generate_retry_fn,
    client,
    progress_cb=None,
    resume_from_batch=0,
    doc_profile=None,
):
    prompt = _build_source_text_prompt(prompt, source_kind=payload_kind, doc_profile=doc_profile)
    clean = clean_text_fn(full_text)
    chunks = [clean[i:i + text_chunk_chars] for i in range(0, len(clean), text_chunk_chars)]
    batches = batch_text_chunks_fn(chunks)
    total_batches = max(1, len(batches))
    start_batch = max(0, min(int(resume_from_batch or 0), total_batches))
    for index, chunk in enumerate(batches[start_batch:], start=start_batch + 1):
        if pause_cb:
            pause_cb()
        payload = chunk + "\n\n" + prompt if ("claude" in model or "gpt" in model) else [chunk, prompt]
        key = build_request_cache_key_fn(model, prompt, payload_kind, sha256_text_fn(chunk))
        stream_with_cache_fn(
            key,
            lambda payload=payload: generate_retry_fn(client, model, payload, log_cb=log_cb),
            out,
            fmt,
            f_obj,
            mem,
            log_cb,
            pause_cb=pause_cb,
        )
        if progress_cb:
            progress_cb(index, total_batches)


def estimate_text_work_units(
    path,
    ext,
    *,
    text_chunk_chars,
    csv_to_accessible_text_fn,
    clean_text_fn,
    batch_text_chunks_fn,
    docx_module,
    openpyxl_module,
    subprocess_module=subprocess,
    fitz_module=None,
):
    if ext in FITZ_TEXT_EXTENSIONS:
        full = _load_fitz_text_content(path, fitz_module=fitz_module)
    else:
        full = _load_text_document_content(
            path,
            ext,
            csv_to_accessible_text_fn=csv_to_accessible_text_fn,
            docx_module=docx_module,
            openpyxl_module=openpyxl_module,
            subprocess_module=subprocess_module,
        )
    clean = clean_text_fn(full)
    chunks = [clean[i:i + text_chunk_chars] for i in range(0, len(clean), text_chunk_chars)]
    batches = batch_text_chunks_fn(chunks)
    return max(1, len(batches))


def process_text(
    client,
    path,
    out,
    ext,
    fmt,
    prompt,
    model,
    f_obj,
    mem,
    log_cb,
    *,
    pause_cb=None,
    text_chunk_chars,
    csv_to_accessible_text_fn,
    clean_text_fn,
    batch_text_chunks_fn,
    build_request_cache_key_fn,
    sha256_text_fn,
    stream_with_cache_fn,
    generate_retry_fn,
    docx_module,
    openpyxl_module,
    subprocess_module=subprocess,
    page_progress_cb=None,
    resume_from_batch=0,
    doc_profile=None,
    fitz_module=None,
):
    if fmt == "html" and doc_profile == "tabular" and ext in TABULAR_TEXT_EXTENSIONS:
        datasets = _build_tabular_datasets(path, ext, openpyxl_module=openpyxl_module)
        if datasets:
            title = os.path.splitext(os.path.basename(path))[0]
            fragment = build_tabular_html_fragment(title, datasets)
            _write_direct_output(fragment, f_obj=f_obj, mem=mem)
            if page_progress_cb:
                page_progress_cb(1, 1)
            return

    if ext in FITZ_TEXT_EXTENSIONS:
        full = _load_fitz_text_content(path, fitz_module=fitz_module)
    else:
        full = _load_text_document_content(
            path,
            ext,
            csv_to_accessible_text_fn=csv_to_accessible_text_fn,
            docx_module=docx_module,
            openpyxl_module=openpyxl_module,
            subprocess_module=subprocess_module,
        )

    return _stream_text_batches(
        full,
        out=out,
        fmt=fmt,
        prompt=prompt,
        model=model,
        f_obj=f_obj,
        mem=mem,
        log_cb=log_cb,
        pause_cb=pause_cb,
        payload_kind="text",
        text_chunk_chars=text_chunk_chars,
        clean_text_fn=clean_text_fn,
        batch_text_chunks_fn=batch_text_chunks_fn,
        build_request_cache_key_fn=build_request_cache_key_fn,
        sha256_text_fn=sha256_text_fn,
        stream_with_cache_fn=stream_with_cache_fn,
        generate_retry_fn=generate_retry_fn,
        client=client,
        progress_cb=page_progress_cb,
        resume_from_batch=resume_from_batch,
        doc_profile=doc_profile,
    )


def process_pptx(
    client,
    path,
    out,
    fmt,
    prompt,
    model,
    f_obj,
    mem,
    log_cb,
    *,
    pause_cb=None,
    page_progress_cb=None,
    text_chunk_chars,
    clean_text_fn,
    batch_text_chunks_fn,
    build_request_cache_key_fn,
    sha256_text_fn,
    stream_with_cache_fn,
    generate_retry_fn,
    subprocess_module=subprocess,
    pptx_module=None,
    resume_from_batch=0,
):
    ext = os.path.splitext(path)[1].lower()
    full = ""
    if ext == ".ppt":
        try:
            conv = subprocess_module.run(
                ["textutil", "-convert", "txt", "-stdout", path],
                capture_output=True,
                text=True,
                check=False,
            )
            if conv.returncode == 0 and conv.stdout.strip():
                full = conv.stdout
            else:
                raise RuntimeError("textutil conversion returned no content.")
        except Exception as ex:
            raise Exception(
                "Legacy .ppt parsing failed. Convert to .pptx for full structural extraction. "
                f"Details: {ex}"
            )
    else:
        if pptx_module is None:
            import pptx as pptx_module  # type: ignore
        prs = pptx_module.Presentation(path)
        total_slides = max(1, len(prs.slides))
        num = 1
        for slide in prs.slides:
            full += f"\n[--- Slide: {num} ---]\n"
            try:
                shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left))
            except Exception:
                shapes = slide.shapes
            for shape in shapes:
                if shape.has_text_frame:
                    full += shape.text + "\n"
            num += 1

    return _stream_text_batches(
        full,
        out=out,
        fmt=fmt,
        prompt=prompt,
        model=model,
        f_obj=f_obj,
        mem=mem,
        log_cb=log_cb,
        pause_cb=pause_cb,
        payload_kind="pptx",
        text_chunk_chars=text_chunk_chars,
        clean_text_fn=clean_text_fn,
        batch_text_chunks_fn=batch_text_chunks_fn,
        build_request_cache_key_fn=build_request_cache_key_fn,
        sha256_text_fn=sha256_text_fn,
        stream_with_cache_fn=stream_with_cache_fn,
        generate_retry_fn=generate_retry_fn,
        client=client,
        progress_cb=page_progress_cb,
        resume_from_batch=resume_from_batch,
        doc_profile="transcript",
    )


def process_epub(
    client,
    path,
    out,
    fmt,
    prompt,
    model,
    f_obj,
    mem,
    log_cb,
    *,
    pause_cb=None,
    page_progress_cb=None,
    text_chunk_chars,
    clean_text_fn,
    batch_text_chunks_fn,
    build_request_cache_key_fn,
    sha256_text_fn,
    stream_with_cache_fn,
    generate_retry_fn,
    epub_module,
    resume_from_batch=0,
):
    book = epub_module.read_epub(path)
    full = ""
    for item in book.get_items():
        if item.get_type() == 9:
            full += re.sub(r"<[^>]+>", " ", item.get_body_content().decode("utf-8")) + "\n"
    return _stream_text_batches(
        full,
        out=out,
        fmt=fmt,
        prompt=prompt,
        model=model,
        f_obj=f_obj,
        mem=mem,
        log_cb=log_cb,
        pause_cb=pause_cb,
        payload_kind="epub",
        text_chunk_chars=text_chunk_chars,
        clean_text_fn=clean_text_fn,
        batch_text_chunks_fn=batch_text_chunks_fn,
        build_request_cache_key_fn=build_request_cache_key_fn,
        sha256_text_fn=sha256_text_fn,
        stream_with_cache_fn=stream_with_cache_fn,
        generate_retry_fn=generate_retry_fn,
        client=client,
        progress_cb=page_progress_cb,
        resume_from_batch=resume_from_batch,
        doc_profile="book",
    )


def process_img(
    client,
    path,
    out,
    fmt,
    prompt,
    model,
    f_obj,
    mem,
    log_cb,
    *,
    pause_cb=None,
    enhance_image_fn,
    build_payload_fn,
    build_request_cache_key_fn,
    sha256_file_fn,
    stream_with_cache_fn,
    generate_retry_fn,
    remove_fn=os.remove,
    image_prepare_fn=prepare_image_for_scan,
):
    if pause_cb:
        pause_cb()
    staged = image_prepare_fn(path, log_cb=log_cb, remove_fn=remove_fn)
    enhanced = enhance_image_fn(staged)
    try:
        if "claude" in model or "gpt" in model:
            payload = build_payload_fn(model, prompt, enhanced, "image/png")
            key = build_request_cache_key_fn(model, prompt, "image", sha256_file_fn(enhanced))
            stream_with_cache_fn(
                key,
                lambda: generate_retry_fn(client, model, payload, log_cb=log_cb),
                out,
                fmt,
                f_obj,
                mem,
                log_cb,
                pause_cb=pause_cb,
            )
        else:
            key = build_request_cache_key_fn(model, prompt, "image-upload", sha256_file_fn(enhanced))

            def _gemini_img_request():
                upload_config = {
                    "mime_type": "image/png",
                    "display_name": os.path.basename(enhanced),
                    "http_options": {"timeout": GEMINI_FILE_UPLOAD_TIMEOUT_MS},
                }
                log_cb(f"[Gemini Image] Uploading image {os.path.basename(enhanced)}.")
                uploaded = client.files.upload(file=enhanced, config=upload_config)
                log_cb(f"[Gemini Image] Waiting for image {os.path.basename(enhanced)} to become ready.")
                uploaded = wait_for_gemini_upload_ready(client, uploaded, poll_sec=0.5, max_wait_sec=30.0, log_cb=log_cb)
                log_cb(f"[Gemini Image] Image {os.path.basename(enhanced)} is ready.")
                response = generate_retry_fn(client, model, [uploaded, prompt], log_cb=log_cb)

                def _cleanup_upload():
                    try:
                        client.files.delete(
                            name=uploaded.name,
                            config={"http_options": {"timeout": GEMINI_FILE_DELETE_TIMEOUT_MS}},
                        )
                    except TypeError:
                        try:
                            client.files.delete(name=uploaded.name)
                        except Exception as cleanup_ex:
                            log_cb(f"Warning: could not delete temporary Gemini upload {uploaded.name}: {cleanup_ex}")
                    except Exception as cleanup_ex:
                        log_cb(f"Warning: could not delete temporary Gemini upload {uploaded.name}: {cleanup_ex}")

                return response, _cleanup_upload

            stream_with_cache_fn(key, _gemini_img_request, out, fmt, f_obj, mem, log_cb, pause_cb=pause_cb)
    finally:
        for cleanup_path in (enhanced, staged):
            if cleanup_path != path:
                try:
                    remove_fn(cleanup_path)
                except FileNotFoundError:
                    pass


def process_rendered_document(
    client,
    path,
    out,
    fmt,
    prompt,
    model,
    f_obj,
    mem,
    log_cb,
    *,
    pause_cb=None,
    page_progress_cb=None,
    fitz_module=None,
    enhance_image_fn,
    build_payload_fn,
    build_request_cache_key_fn,
    sha256_file_fn,
    stream_with_cache_fn,
    generate_retry_fn,
    remove_fn=os.remove,
    mkstemp_fn=tempfile.mkstemp,
    close_fn=os.close,
    resume_from_page=0,
):
    if fitz_module is None:
        import fitz as fitz_module  # type: ignore

    doc = fitz_module.open(path)
    total_pages = max(1, len(doc))
    try:
        start_page = max(0, min(int(resume_from_page or 0), total_pages))
        for page_index in range(start_page, total_pages):
            if pause_cb:
                pause_cb()
            page = doc[page_index]
            matrix = fitz_module.Matrix(2.0, 2.0) if hasattr(fitz_module, "Matrix") else None
            pix = page.get_pixmap(matrix=matrix, alpha=False) if matrix is not None else page.get_pixmap(alpha=False)
            image_bytes = pix.tobytes("png")
            fd, temp_png = mkstemp_fn(prefix="chronicle_rendered_doc_", suffix=".png")
            close_fn(fd)
            try:
                with open(temp_png, "wb") as fh:
                    fh.write(image_bytes)
                page_prompt = (
                    f"{prompt}\n\nRENDERED DOCUMENT PAGE MODE:\n"
                    f"- This is page {page_index + 1} of {total_pages} rendered from {os.path.basename(path)}.\n"
                    "- Read the visible page image faithfully and preserve meaningful reading order."
                )
                process_img(
                    client,
                    temp_png,
                    out,
                    fmt,
                    page_prompt,
                    model,
                    f_obj,
                    mem,
                    log_cb,
                    pause_cb=pause_cb,
                    enhance_image_fn=enhance_image_fn,
                    build_payload_fn=build_payload_fn,
                    build_request_cache_key_fn=build_request_cache_key_fn,
                    sha256_file_fn=sha256_file_fn,
                    stream_with_cache_fn=stream_with_cache_fn,
                    generate_retry_fn=generate_retry_fn,
                    remove_fn=remove_fn,
                    image_prepare_fn=lambda image_path, **_kwargs: image_path,
                )
            finally:
                try:
                    remove_fn(temp_png)
                except FileNotFoundError:
                    pass
            if page_progress_cb:
                page_progress_cb(page_index + 1, total_pages, page_index + 1)
    finally:
        doc.close()
